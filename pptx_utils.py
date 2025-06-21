# pptx_utils.py
# An enhanced utility library to create reliable, non-corrupt presentations.

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

class PptxHelper:
    """
    A robust helper class to simplify the creation of PowerPoint presentations.
    This version is designed to prevent file corruption by using explicit layouts
    and building content slides from a blank canvas.
    """
    def __init__(self):
        self.prs = Presentation()
        # Set a standard 16:9 slide size, common for modern presentations.
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)

    def add_title_slide(self, title, subtitle):
        """Adds a title slide using the guaranteed 'Title Slide' layout (index 0)."""
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        return slide

    def add_blank_slide(self):
        """Adds a blank slide using the guaranteed 'Blank' layout (index 6)."""
        blank_layout = self.prs.slide_layouts[6]
        return self.prs.slides.add_slide(blank_layout)

    def add_textbox(self, slide, text, left, top, width, height, font_size=32, bold=True):
        """Adds a simple textbox with formatting."""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        return txBox

    def add_bullet_list(self, slide, items, left, top, width, height, font_size=18):
        """
        Adds a formatted bullet list to a slide. This is key for building from scratch.
        """
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.clear() # Clear existing paragraph
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)
            p.level = 0 # This creates the bullet point
        return txBox

    def add_image_to_slide(self, slide, image_path, left, top, width):
        """Adds an image with a specified position and size, maintaining aspect ratio."""
        return slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width))

    def add_end_slide(self, title="Thank You", subtitle="Questions?"):
        """Adds a professional, centered final slide on a blank canvas."""
        slide = self.add_blank_slide()
        self.add_textbox(slide, title, left=0, top=1.5, width=10, height=2, font_size=44)
        p = slide.shapes[0].text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        
        self.add_textbox(slide, subtitle, left=0, top=3.0, width=10, height=1.5, font_size=28, bold=False)
        p = slide.shapes[1].text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        return slide

    def save(self, output_path):
        """Saves the presentation to a file."""
        self.prs.save(output_path)